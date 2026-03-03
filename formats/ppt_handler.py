"""
formats/ppt_handler.py — PowerPoint text and table extraction.

Extracts all text runs and tables from a .pptx file and returns them
as a structured DataFrame for downstream query/analysis.

Public API
----------
PPTHandler()
    .extract(path: Path | str) -> pd.DataFrame
      Columns: slide, content_type, content
"""
from __future__ import annotations

import logging
from pathlib import Path

import pandas as pd

logger = logging.getLogger(__name__)


class PPTHandler:
    """Extracts text and tabular content from PowerPoint presentations."""

    def extract(self, path: Path | str) -> pd.DataFrame:
        """
        Parse a .pptx file and return a DataFrame with columns:
        [slide, content_type, content]
        """
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Pt  # type: ignore
        except ImportError as exc:
            raise ImportError("Install python-pptx: pip install python-pptx") from exc

        path = Path(path)
        prs = Presentation(str(path))
        rows: list[dict] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                # ── Text frames ────────────────────────────────────────────
                if shape.has_text_frame:
                    text = "\n".join(
                        para.text for para in shape.text_frame.paragraphs if para.text.strip()
                    )
                    if text.strip():
                        rows.append({
                            "slide": slide_num,
                            "content_type": "text",
                            "content": text.strip(),
                        })

                # ── Tables ─────────────────────────────────────────────────
                if shape.has_table:
                    table = shape.table
                    headers = [cell.text.strip() for cell in table.rows[0].cells]
                    for row_idx, row in enumerate(table.rows[1:], start=1):
                        row_data = {h: row.cells[i].text.strip() for i, h in enumerate(headers)}
                        rows.append({
                            "slide": slide_num,
                            "content_type": "table_row",
                            "content": str(row_data),
                        })

        if not rows:
            logger.warning("No extractable content found in %s", path.name)

        df = pd.DataFrame(rows, columns=["slide", "content_type", "content"])
        logger.info("PPT extracted: %s → %d items", path.name, len(df))
        return df
